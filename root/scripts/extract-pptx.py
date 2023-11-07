#!/bin/python3

from pptx import Presentation
from pptx.text.text import _Paragraph
from pptx.shapes.placeholder import SlidePlaceholder
from pptx.enum.shapes import PP_PLACEHOLDER


def _print_indented_text(paragraph: _Paragraph, bullet: str):
    """
    定义一个函数，用于打印带缩进和bullet的段落文本
    """
    if len(paragraph.text) == 0:
        return
    indent_level = paragraph.level
    # 根据缩进级别添加空格
    indent = '    ' * indent_level
    # 打印缩进和bullet后的段落文本
    print(f"{indent}{bullet} ", end='')
    for run in paragraph.runs:
        if run.font.bold:
            print(f" **{run.text.strip()}** ", end='')
        else:
            print(run.text, end='')
    print()


def extract_pptx(pptx_file: str):
    prs = Presentation(pptx_file)
    for i, slide in enumerate(prs.slides):
        print(f"=== Slide {i + 1}\n")
        for shape in slide.shapes:
            if shape.is_placeholder:
                shape: SlidePlaceholder
                if hasattr(shape, "text") and len(shape.text.strip()) > 0:
                    bullet = '-'
                    match shape.placeholder_format.type:
                        case PP_PLACEHOLDER.TITLE:
                            bullet=''
                            print("##", end="")
                        case PP_PLACEHOLDER.SUBTITLE:
                            bullet=''
                    # 遍历文本框中的每一个段落
                    for paragraph in shape.text_frame.paragraphs:
                        _print_indented_text(paragraph, bullet)
                        pass
                    print()


if __name__ == '__main__':
    from sys import argv
    if len(argv) != 2:
        print("Usage: extract-pptx.py <pptx file>")
        exit(1)
    extract_pptx(argv[1])
